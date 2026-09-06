"""
Project : Converigo
Author  : Converigo Factory (Jalur 2 / F7)
Version : 1.0.0

In-image probe for Factory Batch F7 (office/document cluster:
DOC-14 docx-to-html + pptx-to-png).

Executed INSIDE the production image by docker-runtime-verify step [4/5]
(dispatch with probe_script=scripts/ci_in_document_factory_probe.py):

    python scripts/ci_in_document_factory_probe.py

All fixtures are generated in-image (python-docx / python-pptx), so the
probe is self-sufficient exactly like the F1-F6 probes.  The two F7 net-new
plugins are resolved through the real registry and executed through their
public async convert(); the honest-error path and the office-cluster
contract artifacts (<slug>.contract.json + <slug>.json) are asserted.
Exit code 0 = PASS.
"""
from __future__ import annotations

import asyncio
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.plugins.registry import registry  # noqa: E402

F7_SLUGS = ("docx-to-html", "pptx-to-png")


def _build_docx(path: Path) -> Path:
    from docx import Document

    document = Document()
    document.add_heading("F7 probe document", level=1)
    document.add_paragraph("alpha paragraph")
    document.add_paragraph("beta paragraph")
    document.save(str(path))
    return path


def _build_pptx(path: Path) -> Path:
    from pptx import Presentation

    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]  # title + content layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "F7 probe slide"
    body = slide.placeholders[1].text_frame
    body.text = "alpha bullet"
    bullet = body.add_paragraph()
    bullet.text = "beta bullet"
    presentation.save(str(path))
    return path


def _verify_html(payload: Path) -> None:
    text = payload.read_text(encoding="utf-8")
    assert text.startswith("<!DOCTYPE html>"), text[:64]
    assert "alpha paragraph" in text and "beta paragraph" in text


def _verify_png(payload: Path) -> None:
    from PIL import Image

    assert payload.stat().st_size > 0
    with Image.open(str(payload)) as image:
        assert image.format == "PNG", image.format
        assert image.width > 0 and image.height > 0


async def _convert(slug: str, source: Path, target: str, working: Path) -> Path:
    plugin = registry.by_slug[slug]
    return await plugin.convert(source, target, output_dir=working)


def main() -> int:
    failures: list[str] = []
    with tempfile.TemporaryDirectory(prefix="f7_probe_") as tmp:
        root = Path(tmp)
        docx_fixture = _build_docx(root / "probe_fixture.docx")
        pptx_fixture = _build_pptx(root / "probe_fixture.pptx")

        plan = [
            ("docx-to-html", docx_fixture, "html", _verify_html),
            ("pptx-to-png", pptx_fixture, "png", _verify_png),
        ]
        for slug, fixture, target, verifier in plan:
            try:
                assert registry.has_slug(slug), f"{slug} not registered"
                output_path = asyncio.run(_convert(
                    slug, fixture, target, root / f"out_{slug.replace('-', '_')}"
                ))
                assert output_path.is_file() and output_path.stat().st_size > 0
                verifier(output_path)
                print(f"F7 PROBE OK: {slug} ({fixture.suffix} -> {target})")
            except Exception as exc:  # noqa: BLE001 - probe reports all
                failures.append(f"{slug}: {type(exc).__name__}: {exc}")

        # Honest error: a plain-text file must never fabricate a conversion.
        try:
            bogus = root / "probe_bogus.txt"
            bogus.write_text("definitely not a document", encoding="utf-8")
            asyncio.run(_convert("docx-to-html", bogus, "html", root / "out_bogus"))
            failures.append("docx-to-html: bogus input did not fail honestly")
        except Exception:
            print("F7 PROBE OK: honest error for non-DOCX input")

        converters_dir = (
            Path(__file__).resolve().parent.parent / "app" / "data" / "converters"
        )
        for slug in F7_SLUGS:
            for suffix in ("contract.json", "json"):
                artifact = converters_dir / f"{slug}.{suffix}"
                if artifact.exists():
                    print(f"F7 PROBE OK: artifact {artifact.name}")
                else:
                    failures.append(f"artifact missing: {artifact.name}")

    if failures:
        print("F7 PROBE: FAIL")
        for item in failures:
            print(f"  - {item}")
        return 1
    print("F7 PROBE: PASS (2/2 document converters verified in-image)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())