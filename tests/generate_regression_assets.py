from PIL import Image, ImageDraw, ImageFont
from pathlib import Path
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

OUT = Path('tests/assets/regression')
OUT.mkdir(parents=True, exist_ok=True)

def make_image(path, size=(800,600), color=(200,60,80), text=None, fmt='JPEG'):
    img = Image.new('RGB', size, color)
    d = ImageDraw.Draw(img)
    if text:
        try:
            f = ImageFont.load_default()
            d.text((10,10), text, fill=(255,255,255), font=f)
        except Exception:
            d.text((10,10), text, fill=(255,255,255))
    img.save(path, fmt)
    print('WROTE:', path, path.stat().st_size)

def make_pdf(path):
    img_path = OUT / 'tmp_pdf_page.png'
    make_image(img_path, text='PDF page', fmt='PNG')
    img = Image.open(img_path)
    img.save(path, 'PDF', resolution=100.0)
    img_path.unlink()
    print('WROTE:', path, path.stat().st_size)

def make_docx(path):
    doc = Document()
    doc.add_heading('Sample DOCX', level=1)
    doc.add_paragraph('This is a generated test document for regression.')
    doc.save(path)
    print('WROTE:', path, path.stat().st_size)

def make_xlsx(path):
    wb = Workbook()
    ws = wb.active
    ws.title = 'Sheet1'
    ws['A1'] = 'Sample'
    ws['B1'] = 123
    wb.save(path)
    print('WROTE:', path, path.stat().st_size)

def make_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = top = Inches(1)
    width = Inches(8)
    height = Inches(4.5)
    tx = slide.shapes.add_textbox(left, top, width, height).text_frame
    tx.text = 'Generated test PPTX slide.'
    prs.save(path)
    print('WROTE:', path, path.stat().st_size)

def make_mp4(path, duration=2, size=(320,240), color=(0,128,200)):
    try:
        from moviepy.video.VideoClip import ColorClip
    except Exception as e:
        print('moviepy not available, skipping mp4:', e)
        return
    clip = ColorClip(size, color=color, duration=duration)
    try:
        clip = clip.set_fps(24)
    except Exception:
        try:
            clip.fps = 24
        except Exception:
            pass
    try:
        clip.write_videofile(str(path), codec='libx264', audio=False, fps=getattr(clip, 'fps', 24))
    except Exception as e:
        print('MP4 write failed:', e)
        return
    print('WROTE:', path, path.stat().st_size)

if __name__ == '__main__':
    make_image(OUT / 'sample.jpg', text='sample.jpg', fmt='JPEG')
    make_image(OUT / 'sample.png', text='sample.png', fmt='PNG')
    try:
        make_image(OUT / 'sample.webp', text='sample.webp', fmt='WEBP')
    except Exception as e:
        print('WEBP skipped:', e)
    make_pdf(OUT / 'sample.pdf')
    make_docx(OUT / 'sample.docx')
    make_xlsx(OUT / 'sample.xlsx')
    make_pptx(OUT / 'sample.pptx')
    try:
        make_mp4(OUT / 'sample.mp4')
    except Exception as e:
        print('MP4 skipped:', e)
