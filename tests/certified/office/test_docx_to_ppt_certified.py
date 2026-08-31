from pathlib import Path

from fastapi.testclient import TestClient
from pptx import Presentation

from app.core.settings import settings
from app.main import app


def _all_slide_text(presentation: Presentation) -> str:
    """Concatenate every text run in the presentation (titles, bodies, tables)."""
    chunks = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        chunks.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            chunks.append(cell.text.strip())
    return "\n".join(chunks)


def test_docx_to_ppt_conversion_creates_valid_pptx(tmp_path: Path):
    client = TestClient(app)

    docx_path = Path("tests/assets/regression/sample.docx")
    assert docx_path.exists(), "Sample DOCX is missing"

    resp = client.post(
        "/convert",
        data={"target_format": "pptx", "operation": "docx-to-ppt"},
        files={
            "file": (
                docx_path.name,
                docx_path.read_bytes(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )

    assert resp.status_code == 201, resp.text
    payload = resp.json()
    assert payload.get("status") == "success"
    download_path = payload.get("download_path")
    assert download_path, payload
    assert download_path.startswith("/download/")
    relative_parts = Path(download_path.removeprefix("/download/")).parts
    assert len(relative_parts) == 2, f"Unexpected download path shape: {download_path}"
    conversion_id, filename = relative_parts
    assert filename.endswith(".pptx")

    local_path = settings.OUTPUT_DIR / conversion_id / filename
    download_resp = client.get(download_path)
    assert download_resp.status_code == 200, download_resp.text
    assert download_resp.content, "Downloaded content is empty"

    assert local_path.exists(), f"Expected output PPTX not found: {local_path}"
    assert local_path.stat().st_size > 0, "Output PPTX is empty"
    assert local_path.suffix.lower() == ".pptx"

    # Content-integrity: re-open output with python-pptx and verify the source
    # DOCX text actually made it into the presentation (not a fake file).
    presentation = Presentation(str(local_path))
    slides = list(presentation.slides)
    assert len(slides) >= 1, "PPTX output has no slides"
    joined = _all_slide_text(presentation)
    assert "Sample DOCX" in joined, f"Expected source paragraph in PPTX, got: {joined[:500]}"
    assert "This is a generated test document for regression." in joined, (
        f"Expected source paragraph in PPTX, got: {joined[:500]}"
    )

    assert len(list((settings.OUTPUT_DIR / conversion_id).glob("*"))) == 1


def test_docx_to_ppt_output_is_real_ooxml_pptx(tmp_path: Path):
    """Verify the output starts with the OOXML ZIP magic bytes (PK), proving it
    is a real PPTX container rather than a plain-text/fake file."""
    client = TestClient(app)

    docx_path = Path("tests/assets/regression/sample.docx")
    resp = client.post(
        "/convert",
        data={"target_format": "pptx", "operation": "docx-to-ppt"},
        files={
            "file": (
                docx_path.name,
                docx_path.read_bytes(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )
    assert resp.status_code == 201, resp.text
    payload = resp.json()
    conversion_id, filename = Path(payload["download_path"].removeprefix("/download/")).parts
    local_path = settings.OUTPUT_DIR / conversion_id / filename

    with local_path.open("rb") as handle:
        header = handle.read(4)
    assert header == b"PK\x03\x04", f"Output is not a real PPTX (ZIP) file: {header!r}"
