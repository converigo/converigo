"""
Project : Converigo
Author  : Pico Lala & ChatGPT
Version : 3.0.1

DOCX -> PPT Plugin

Converts Word documents into PowerPoint presentations. Source detection is
content-based (magic bytes) rather than extension-based:

- ZIP/OOXML magic (PK) -> treated as DOCX and real document content is extracted,
  which also covers `.doc` files that are actually DOCX renamed.
- OLE2/CFB magic (d0cf11e0a1b11ae1) -> legacy binary `.doc`; not supported,
  raises an explicit error instead of emitting an empty PPTX.
- Anything else -> raises an explicit error (no silent fake pass).

Mapping (v1):
- One slide per top-level heading; non-heading paragraphs are appended to the
  current slide body until MAX_PARAGRAPH_CHARS_PER_SLIDE, then a new slide starts.
- Native Word tables are mapped to native python-pptx tables when the column
  count is <= MAX_TABLE_NATIVE_COLUMNS; wider tables are rendered as text.
- Inline pictures are copied onto a slide when available (best-effort).
"""

import logging

from pathlib import Path

from app.plugins.base import ConverterPlugin

logger = logging.getLogger(__name__)

# Tunable mapping flags (see plan: .tmp/PLAN_PR1C_DOCX_PPT_XLSX_PPT.md).
MAX_TABLE_NATIVE_COLUMNS = 8
MAX_PARAGRAPH_CHARS_PER_SLIDE = 2000
MAX_ROWS_PER_SLIDE = 20


class DOCXToPPTPlugin(ConverterPlugin):
    slug = "docx-to-ppt"
    name = "DOCX to PPT"
    description = "Convert Word documents into PowerPoint presentations."
    category = "document"
    engine = "document"
    icon = "📄"

    source_formats = ["docx", "doc", "word"]
    target_formats = ["ppt", "pptx", "powerpoint"]

    goal = "document"
    use_case = "Best for turning Word documents into editable PowerPoint presentations."
    priority = 80
    quality = 85
    compatibility = 80
    estimated_saving = 10
    badge = "Office Conversion"
    seo_title = "DOCX to PPT Converter | Converigo"
    seo_description = "Convert Word documents into PowerPoint presentations quickly and easily."

    # DOCX (OOXML) files are ZIP archives and start with the "PK" signature.
    _DOCX_PK_MAGIC = b"PK\x03\x04"
    # Legacy .doc files are OLE2 / Compound File Binary containers.
    _DOC_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

    @classmethod
    def _detect_container(cls, source_path: Path) -> str:
        """Return the container type detected from file content.

        Returns "docx" for ZIP/OOXML files, "doc" for OLE2/CFB files, or
        "unknown" for anything else (empty file, RTF, plain text, ...).
        """
        try:
            with source_path.open("rb") as fh:
                header = fh.read(8)
        except OSError as exc:
            raise RuntimeError(f"Could not read the source document: {exc}") from exc

        if header.startswith(cls._DOCX_PK_MAGIC):
            return "docx"
        if header.startswith(cls._DOC_OLE2_MAGIC):
            return "doc"
        return "unknown"

    @staticmethod
    def _is_heading(paragraph) -> bool:
        """Return True when the paragraph uses a built-in heading style."""
        style_name = getattr(paragraph.style, "name", "") or ""
        return style_name.lower().startswith("heading")

    @staticmethod
    def _new_slide(presentation, title_text: str):
        """Add a slide and set its title (falling back to a text box)."""
        if len(presentation.slide_layouts) > 1:
            layout = presentation.slide_layouts[1]  # Title and Content
        else:
            layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        else:
            from pptx.util import Inches

            title_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.3),
                Inches(9.0),
                Inches(0.8),
            )
            title_box.text_frame.text = title_text
        return slide

    @staticmethod
    def _append_body_text(slide, text: str) -> None:
        """Append a paragraph of text to a slide body text frame."""
        if slide.shapes.placeholders and len(slide.shapes.placeholders) > 1:
            body = slide.shapes.placeholders[1]
            text_frame = body.text_frame
        else:
            from pptx.util import Inches

            textbox = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(1.2),
                Inches(9.0),
                Inches(6.0),
            )
            text_frame = textbox.text_frame

        if text_frame.paragraphs and not text_frame.paragraphs[0].runs:
            text_frame.paragraphs[0].text = text
        else:
            text_frame.add_paragraph().text = text

    @classmethod
    def _append_native_table(cls, slide, docx_table) -> None:
        """Copy a Word table into the current slide as a native pptx table."""
        rows = len(docx_table.rows)
        cols = len(docx_table.columns)
        if not rows or not cols:
            return

        from pptx.util import Inches

        graphic_frame = slide.shapes.add_table(
            rows,
            cols,
            Inches(0.5),
            Inches(1.2),
            Inches(9.0),
            Inches(0.3 * min(rows, MAX_ROWS_PER_SLIDE) + 0.5),
        )
        table = graphic_frame.table
        for r in range(rows):
            for c in range(cols):
                table.cell(r, c).text = docx_table.cell(r, c).text.strip()

    async def convert(
        self,
        source_path: Path,
        target_format: str,
        output_dir: Path | None = None,
        temp_dir: Path | None = None,
    ) -> Path:
        if not self.supports(source_path.suffix, target_format):
            raise RuntimeError("DOCXToPPTPlugin only supports DOCX/DOC -> PPT.")

        container = self._detect_container(source_path)

        if container == "doc":
            raise RuntimeError(
                "Legacy .doc format (OLE2 compound document) is not supported by "
                "this converter. Please save your document as .docx and try again."
            )
        if container == "unknown":
            raise RuntimeError(
                "The uploaded file is not a valid DOCX document. Please save it "
                "as a valid .docx file and try again."
            )

        try:
            from docx import Document

            document = Document(str(source_path))
        except Exception as exc:
            raise RuntimeError(
                "The DOCX document could not be parsed. Please save it as a valid "
                ".docx file and try again."
            ) from exc

        from pptx import Presentation

        presentation = Presentation()

        paragraphs = [p for p in document.paragraphs if p.text.strip()]
        slide = None
        chars_on_slide = 0

        if not paragraphs:
            # Empty document: emit a single honest slide so the output is a
            # valid (non-fake) PPTX explaining there was nothing to convert.
            slide = self._new_slide(presentation, source_path.stem)
            self._append_body_text(slide, "This document contains no text content.")
        else:
            for paragraph in paragraphs:
                text = paragraph.text.strip()
                is_heading = self._is_heading(paragraph)

                if (
                    slide is None
                    or is_heading
                    or chars_on_slide >= MAX_PARAGRAPH_CHARS_PER_SLIDE
                ):
                    slide = self._new_slide(
                        presentation,
                        text if is_heading else source_path.stem,
                    )
                    chars_on_slide = 0
                    if is_heading:
                        continue

                self._append_body_text(slide, text)
                chars_on_slide += len(text)

        # Native tables: one slide per table (generic label when no caption).
        for index, docx_table in enumerate(document.tables, start=1):
            if not docx_table.rows:
                continue
            cols = len(docx_table.columns)
            table_slide = self._new_slide(presentation, f"Table {index}")
            if cols <= MAX_TABLE_NATIVE_COLUMNS:
                self._append_native_table(table_slide, docx_table)
            else:
                lines = [
                    " | ".join(cell.text.strip() for cell in row.cells)
                    for row in docx_table.rows
                ]
                for line in lines:
                    self._append_body_text(table_slide, line)

        # Copy the first inline picture onto the last slide when present.
        try:
            from docx.enum.shape import WD_INLINE_SHAPE

            import io

            from pptx.util import Inches

            for inline in document.inline_shapes:
                if inline.type != WD_INLINE_SHAPE.PICTURE:
                    continue
                try:
                    # python-docx 1.1.x InlineShape has no `.blob`/`.image`; read the
                    # image bytes through its relationship instead.
                    blip = inline._inline.graphic.graphicData.pic.blipFill.blip
                    embed_rid = blip.embed
                    blob = document.part.related_parts[embed_rid].blob
                except (AttributeError, KeyError) as exc:
                    logger.warning("Could not extract inline picture from DOCX: %s", exc)
                    blob = None
                if blob:
                    target_slide = (
                        presentation.slides[-1]
                        if presentation.slides
                        else self._new_slide(presentation, source_path.stem)
                    )
                    try:
                        target_slide.shapes.add_picture(
                            io.BytesIO(blob),
                            Inches(0.5),
                            Inches(4.0),
                            width=Inches(4.0),
                        )
                    except (OSError, ValueError) as exc:
                        logger.warning("Could not embed picture into PPTX: %s", exc)
                    else:
                        break
        except Exception as exc:
            # Images are best-effort; text fidelity is the guaranteed core. Log the
            # failure instead of swallowing it silently so regressions stay visible.
            logger.warning("Inline-picture copy step failed: %s", exc)

        from app.core.settings import settings

        working_dir = temp_dir or output_dir or (settings.OUTPUT_DIR / "document")
        working_dir.mkdir(parents=True, exist_ok=True)

        output_path = working_dir / f"{source_path.stem}.pptx"
        presentation.save(str(output_path))

        if not output_path.exists() or output_path.stat().st_size == 0:
            raise RuntimeError("DOCX/DOC to PPT conversion did not produce output.")

        return output_path

