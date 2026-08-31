"""
Project : Converigo
Author  : Pico Lala & ChatGPT
Version : 3.0.1

PPTX -> XLSX Plugin

Converts PowerPoint presentations into Excel spreadsheets. Source detection is
content-based (magic bytes) rather than extension-based:

- ZIP/OOXML magic (PK) -> treated as PPTX and real slide content is extracted,
  which also covers `.ppt` files that are actually PPTX renamed.
- OLE2/CFB magic (d0cf11e0a1b11ae1) -> legacy binary `.ppt`; not supported,
  raises an explicit error instead of emitting an empty XLSX.
- Anything else -> raises an explicit error (no silent fake pass).
"""

from pathlib import Path

from app.plugins.base import ConverterPlugin


class PPTToXLSXPlugin(ConverterPlugin):
    slug = "ppt-to-xlsx"
    name = "PPT to XLSX"
    description = "Convert PowerPoint presentations into Excel spreadsheets."
    category = "document"
    engine = "document"
    icon = "📄"

    source_formats = ["ppt", "pptx", "powerpoint"]
    target_formats = ["xlsx", "xls", "spreadsheet"]

    goal = "document"
    use_case = "Best for extracting PowerPoint slide text and tables into Excel spreadsheets."
    priority = 80
    quality = 85
    compatibility = 80
    estimated_saving = 10
    badge = "Office Conversion"
    seo_title = "PPT to XLSX Converter | Converigo"
    seo_description = "Convert PowerPoint presentations into Excel spreadsheets quickly and easily."

    # PPTX (OOXML) files are ZIP archives and start with the "PK" signature.
    _PPTX_PK_MAGIC = b"PK\x03\x04"
    # Legacy .ppt files are OLE2 / Compound File Binary containers.
    _PPT_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

    @classmethod
    def _detect_container(cls, source_path: Path) -> str:
        """Return the container type detected from file content.

        Returns "pptx" for ZIP/OOXML files, "ppt" for OLE2/CFB files, or
        "unknown" for anything else (empty file, RTF, plain text, ...).
        """
        try:
            with source_path.open("rb") as fh:
                header = fh.read(8)
        except OSError as exc:
            raise RuntimeError(f"Could not read the source presentation: {exc}") from exc

        if header.startswith(cls._PPTX_PK_MAGIC):
            return "pptx"
        if header.startswith(cls._PPT_OLE2_MAGIC):
            return "ppt"
        return "unknown"

    @staticmethod
    def _shape_text(shape) -> str:
        """Return the concatenated, non-empty paragraph text of a shape.

        Handles shapes that expose a text frame (placeholders, text boxes).
        Grouped/nested shapes are out of scope for v1 and are skipped.
        """
        if not getattr(shape, "has_text_frame", False):
            return ""
        lines = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                lines.append(text)
        return "\n".join(lines)

    @staticmethod
    def _append_slide_table(sheet, table_shape) -> None:
        """Copy a slide table into the current worksheet row by row."""
        if not getattr(table_shape, "rows", None):
            return
        for row in table_shape.rows:
            sheet.append([cell.text.strip() for cell in row.cells])

    async def convert(
        self,
        source_path: Path,
        target_format: str,
        output_dir: Path | None = None,
        temp_dir: Path | None = None,
    ) -> Path:
        if not self.supports(source_path.suffix, target_format):
            raise RuntimeError("PPTToXLSXPlugin only supports PPT/PPTX -> XLSX.")

        container = self._detect_container(source_path)

        if container == "ppt":
            raise RuntimeError(
                "Legacy .ppt format (OLE2 compound document) is not supported by "
                "this converter. Please save your presentation as .pptx and try again."
            )
        if container == "unknown":
            raise RuntimeError(
                "The uploaded file is not a valid PPTX presentation. Please save it "
                "as a valid .pptx file and try again."
            )

        try:
            from pptx import Presentation

            presentation = Presentation(str(source_path))
        except Exception as exc:
            raise RuntimeError(
                "The PPTX presentation could not be parsed. Please save it as a valid "
                ".pptx file and try again."
            ) from exc

        import openpyxl

        workbook = openpyxl.Workbook()
        workbook.remove(workbook.active)

        # One sheet per slide: "Slide 1", "Slide 2", ... Each sheet holds the
        # slide text lines followed by any native slide tables (cell layout).
        for index, slide in enumerate(presentation.slides, start=1):
            sheet = workbook.create_sheet(f"Slide {index}")
            sheet.append(["Slide Text"])
            for shape in slide.shapes:
                text = self._shape_text(shape)
                if text:
                    for line in text.split("\n"):
                        sheet.append([line])
                if getattr(shape, "has_table", False):
                    self._append_slide_table(sheet, shape.table)

        from app.core.settings import settings

        working_dir = temp_dir or output_dir or (settings.OUTPUT_DIR / "document")
        working_dir.mkdir(parents=True, exist_ok=True)

        output_path = working_dir / f"{source_path.stem}.xlsx"
        workbook.save(str(output_path))
        workbook.close()

        if not output_path.exists() or output_path.stat().st_size == 0:
            raise RuntimeError("PPT/PPTX to XLSX conversion did not produce output.")

        return output_path
