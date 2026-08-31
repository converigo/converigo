"""
Project : Converigo
Author  : Pico Lala & ChatGPT
Version : 3.0.1

PPTX -> DOCX Plugin

Converts PowerPoint presentations into Word documents. Source detection is
content-based (magic bytes) rather than extension-based:

- ZIP/OOXML magic (PK) -> treated as PPTX and real slide text is extracted,
  which also covers `.ppt` files that are actually PPTX renamed.
- OLE2/CFB magic (d0cf11e0a1b11ae1) -> legacy binary `.ppt`; not supported,
  raises an explicit error instead of emitting a content-less DOCX.
- Anything else -> raises an explicit error (no silent fake pass).
"""

from pathlib import Path

from app.plugins.base import ConverterPlugin


class PPTToDOCXPlugin(ConverterPlugin):
    slug = "ppt-to-docx"
    name = "PPT to DOCX"
    description = "Convert PowerPoint presentations into editable Word documents."
    category = "document"
    engine = "document"
    icon = "📄"

    source_formats = ["ppt", "pptx", "powerpoint"]
    target_formats = ["docx", "doc", "word"]

    goal = "document"
    use_case = "Best for turning PowerPoint slides into editable Word documents."
    priority = 80
    quality = 85
    compatibility = 80
    estimated_saving = 10
    badge = "Office Conversion"
    seo_title = "PPT to DOCX Converter | Converigo"
    seo_description = "Convert PowerPoint presentations into editable Word documents quickly and easily."

    # PPTX (OOXML) files are ZIP archives and start with the "PK" signature.
    _PPTX_PK_MAGIC = b"PK\x03\x04"
    # Legacy .ppt files are OLE2 / Compound File Binary containers.
    _PPT_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

    @classmethod
    def _detect_container(cls, source_path: Path) -> str:
        """Return the container type detected from file content.

        Returns "pptx" for ZIP/OOXML files, "ppt" for OLE2/CFB files, or
        "unknown" for anything else (empty file, RTF, plain text, ...).
        """
        try:
            with source_path.open("rb") as fh:
                header = fh.read(8)
        except OSError as exc:
            raise RuntimeError(f"Could not read the source presentation: {exc}") from exc

        if header.startswith(cls._PPTX_PK_MAGIC):
            return "pptx"
        if header.startswith(cls._PPT_OLE2_MAGIC):
            return "ppt"
        return "unknown"

    @staticmethod
    def _shape_text(shape) -> str:
        """Return the concatenated, non-empty paragraph text of a shape.

        Handles shapes that expose a text frame (placeholders, text boxes).
        Grouped/nested shapes are out of scope for v1 and are skipped.
        """
        if not getattr(shape, "has_text_frame", False):
            return ""
        lines = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                lines.append(text)
        return "\n".join(lines)

    @staticmethod
    def _append_slide_table(document, table_shape) -> None:
        """Copy a slide table into the Word document as a native table."""
        if not getattr(table_shape, "rows", None):
            return
        rows = len(table_shape.rows)
        cols = len(table_shape.columns)
        if not rows or not cols:
            return

        table = document.add_table(rows=rows, cols=cols)
        for r in range(rows):
            for c in range(cols):
                table.cell(r, c).text = table_shape.cell(r, c).text.strip()
        try:
            table.style = "Table Grid"
        except Exception:
            # Style is cosmetic only; a valid table without it is still fine.
            pass

    async def convert(
        self,
        source_path: Path,
        target_format: str,
        output_dir: Path | None = None,
        temp_dir: Path | None = None,
    ) -> Path:
        if not self.supports(source_path.suffix, target_format):
            raise RuntimeError("PPTToDOCXPlugin only supports PPT/PPTX -> DOCX.")

        container = self._detect_container(source_path)

        if container == "ppt":
            raise RuntimeError(
                "Legacy .ppt format (OLE2 compound document) is not supported by "
                "this converter. Please save your presentation as .pptx and try again."
            )
        if container == "unknown":
            raise RuntimeError(
                "The uploaded file is not a valid PPTX presentation. Please save it "
                "as a valid .pptx file and try again."
            )

        try:
            from pptx import Presentation

            presentation = Presentation(str(source_path))
        except Exception as exc:
            raise RuntimeError(
                "The PPTX presentation could not be parsed. Please save it as a valid "
                ".pptx file and try again."
            ) from exc

        from docx import Document

        document = Document()

        for idx, slide in enumerate(presentation.slides, start=1):
            document.add_heading(f"Slide {idx}", level=1)
            for shape in slide.shapes:
                text = self._shape_text(shape)
                if text:
                    document.add_paragraph(text)
                if getattr(shape, "has_table", False):
                    self._append_slide_table(document, shape.table)

        from app.core.settings import settings

        working_dir = temp_dir or output_dir or (settings.OUTPUT_DIR / "document")
        working_dir.mkdir(parents=True, exist_ok=True)

        output_path = working_dir / f"{source_path.stem}.docx"
        document.save(str(output_path))

        if not output_path.exists() or output_path.stat().st_size == 0:
            raise RuntimeError("PPT/PPTX to DOCX conversion did not produce output.")

        return output_path
