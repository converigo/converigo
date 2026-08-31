"""
Project : Converigo
Author  : Pico Lala & ChatGPT
Version : 3.0.1

XLSX -> PPT Plugin

Converts Excel spreadsheets into PowerPoint presentations. Source detection is
content-based (magic bytes) rather than extension-based:

- ZIP/OOXML magic (PK) -> treated as XLSX and real cell data is extracted,
  which also covers `.xls` files that are actually XLSX renamed.
- OLE2/CFB magic (d0cf11e0a1b11ae1) -> legacy binary `.xls`; not supported,
  raises an explicit error instead of emitting a content-less PPTX.
- Anything else -> raises an explicit error (no silent fake pass).

Mapping (v1):
- One summary slide per sheet (sheet name + header row).
- Data rows are mapped to native python-pptx tables when the column count is
  <= MAX_TABLE_NATIVE_COLUMNS; wider sheets are rendered as text lines.
- Large sheets are chunked across slides (MAX_ROWS_PER_SLIDE per slide).
"""

from pathlib import Path

from app.plugins.base import ConverterPlugin

# Tunable mapping flags (see plan: .tmp/PLAN_PR1C_DOCX_PPT_XLSX_PPT.md).
MAX_TABLE_NATIVE_COLUMNS = 8
MAX_ROWS_PER_SLIDE = 20


class XLSXToPPTPlugin(ConverterPlugin):
    slug = "xlsx-to-ppt"
    name = "XLSX to PPT"
    description = "Convert Excel spreadsheets into PowerPoint presentations."
    category = "document"
    engine = "document"
    icon = "📄"

    source_formats = ["xlsx", "xls", "spreadsheet"]
    target_formats = ["ppt", "pptx", "powerpoint"]

    goal = "document"
    use_case = "Best for turning spreadsheet data into PowerPoint presentations."
    priority = 80
    quality = 85
    compatibility = 80
    estimated_saving = 10
    badge = "Office Conversion"
    seo_title = "XLSX to PPT Converter | Converigo"
    seo_description = "Convert Excel spreadsheets into PowerPoint presentations quickly and easily."

    # XLSX (OOXML) files are ZIP archives and start with the "PK" signature.
    _XLSX_PK_MAGIC = b"PK\x03\x04"
    # Legacy .xls files are OLE2 / Compound File Binary containers.
    _XLS_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

    @classmethod
    def _detect_container(cls, source_path: Path) -> str:
        """Return the container type detected from file content.

        Returns "xlsx" for ZIP/OOXML files, "xls" for OLE2/CFB files, or
        "unknown" for anything else (empty file, RTF, plain text, ...).
        """
        try:
            with source_path.open("rb") as fh:
                header = fh.read(8)
        except OSError as exc:
            raise RuntimeError(f"Could not read the source spreadsheet: {exc}") from exc

        if header.startswith(cls._XLSX_PK_MAGIC):
            return "xlsx"
        if header.startswith(cls._XLS_OLE2_MAGIC):
            return "xls"
        return "unknown"

    @staticmethod
    def _cell_text(value) -> str:
        """Render a cell value as plain text for a slide table cell."""
        if value is None:
            return ""
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value)

    @staticmethod
    def _new_slide(presentation, title_text: str):
        """Add a slide and set its title (falling back to a text box)."""
        if len(presentation.slide_layouts) > 1:
            layout = presentation.slide_layouts[1]  # Title and Content
        else:
            layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        else:
            from pptx.util import Inches

            title_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.3),
                Inches(9.0),
                Inches(0.8),
            )
            title_box.text_frame.text = title_text
        return slide

    @staticmethod
    def _append_body_text(slide, text: str) -> None:
        """Append a paragraph of text to a slide body text frame."""
        if slide.shapes.placeholders and len(slide.shapes.placeholders) > 1:
            body = slide.shapes.placeholders[1]
            text_frame = body.text_frame
        else:
            from pptx.util import Inches

            textbox = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(1.2),
                Inches(9.0),
                Inches(6.0),
            )
            text_frame = textbox.text_frame

        if text_frame.paragraphs and not text_frame.paragraphs[0].runs:
            text_frame.paragraphs[0].text = text
        else:
            text_frame.add_paragraph().text = text

    @staticmethod
    def _append_native_rows(slide, rows: list[list[str]]) -> None:
        """Render data rows as a native pptx table on the slide."""
        if not rows:
            return
        rows_count = len(rows)
        cols_count = max(len(row) for row in rows)

        from pptx.util import Inches

        graphic_frame = slide.shapes.add_table(
            rows_count,
            cols_count,
            Inches(0.5),
            Inches(1.2),
            Inches(9.0),
            Inches(0.3 * rows_count + 0.5),
        )
        table = graphic_frame.table
        for r, row in enumerate(rows):
            for c in range(cols_count):
                value = row[c] if c < len(row) else ""
                table.cell(r, c).text = XLSXToPPTPlugin._cell_text(value)

    async def convert(
        self,
        source_path: Path,
        target_format: str,
        output_dir: Path | None = None,
        temp_dir: Path | None = None,
    ) -> Path:
        if not self.supports(source_path.suffix, target_format):
            raise RuntimeError("XLSXToPPTPlugin only supports XLSX/XLS -> PPT.")

        container = self._detect_container(source_path)

        if container == "xls":
            raise RuntimeError(
                "Legacy .xls format (OLE2 compound document) is not supported by "
                "this converter. Please save your spreadsheet as .xlsx and try again."
            )
        if container == "unknown":
            raise RuntimeError(
                "The uploaded file is not a valid XLSX spreadsheet. Please save it "
                "as a valid .xlsx file and try again."
            )

        try:
            import openpyxl

            workbook = openpyxl.load_workbook(
                str(source_path),
                data_only=True,
                read_only=True,
            )
        except Exception as exc:
            raise RuntimeError(
                "The XLSX spreadsheet could not be parsed. Please save it as a valid "
                ".xlsx file and try again."
            ) from exc

        from pptx import Presentation

        presentation = Presentation()

        try:
            for worksheet in workbook.worksheets:
                rows = [
                    row
                    for row in worksheet.iter_rows(values_only=True)
                    if any(value is not None for value in row)
                ]
                if not rows:
                    continue

                cols = max(len(row) for row in rows)
                title = worksheet.title

                if cols <= MAX_TABLE_NATIVE_COLUMNS:
                    for start in range(0, len(rows), MAX_ROWS_PER_SLIDE):
                        chunk = rows[start : start + MAX_ROWS_PER_SLIDE]
                        slide_title = title if start == 0 else f"{title} (continued)"
                        slide = self._new_slide(presentation, slide_title)
                        self._append_native_rows(slide, chunk)
                else:
                    slide = self._new_slide(presentation, title)
                    for row in rows:
                        line = " | ".join(
                            self._cell_text(value) for value in row
                        )
                        if line.strip():
                            self._append_body_text(slide, line)
        finally:
            workbook.close()

        if not presentation.slides:
            slide = self._new_slide(presentation, source_path.stem)
            self._append_body_text(slide, "This spreadsheet contains no data.")

        from app.core.settings import settings

        working_dir = temp_dir or output_dir or (settings.OUTPUT_DIR / "document")
        working_dir.mkdir(parents=True, exist_ok=True)

        output_path = working_dir / f"{source_path.stem}.pptx"
        presentation.save(str(output_path))

        if not output_path.exists() or output_path.stat().st_size == 0:
            raise RuntimeError("XLSX/XLS to PPT conversion did not produce output.")

        return output_path

