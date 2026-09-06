"""
Project : Converigo
Author  : Converigo Factory (Jalur 2 / F7)
Version : 1.0.0

Document Converters (DOC-14 docx-to-html, pptx-to-png)
Factory Batch F7 - office/document cluster, net-new pair.

Built on the F0 certified factory scaffolding: the conversion pipeline
(discovery -> supports() check -> working root -> single servable file ->
non-empty output -> honest error) is owned by FactoryConversionPlugin.
Each converter below is pure configuration plus a small hook.

Semantics (fixed, deterministic, single-servable-file):
- docx-to-html: mammoth.convert_to_html -> semantic HTML fragment (mammoth
  inlines images as base64 data URIs, so the output stays one file), wrapped
  in a minimal deterministic HTML5 document.  Only real OOXML (.docx ZIP)
  containers are accepted; legacy OLE2 .doc and anything else is rejected
  honestly (422 UNSUPPORTED_CONVERSION), never fabricated.
- pptx-to-png: first slide only.  python-pptx extracts the slide's text with
  the exact extraction rules of the certified document engine
  (pptx-to-jpg path), reportlab renders a one-page PDF, PyMuPDF rasterizes
  page 1 to PNG with the engine's rendering profile (fitz.Matrix(2, 2),
  alpha=False).  Presentations without text render the engine's placeholder
  line, mirroring the certified engine behavior.
"""
from __future__ import annotations

import html as _html
from pathlib import Path

from app.factory import make_plugin_class


def _unsupported(source: str, target: str, message: str) -> Exception:
    """Lazily build the honest-422 error (rar-extract/F4 lazy-import precedent)."""
    from app.services.conversion_service import UnsupportedConversionError

    return UnsupportedConversionError(source, target, message)


# OOXML containers (docx/pptx) are ZIP archives starting with the "PK" magic;
# legacy Office binary formats are OLE2 / Compound File Binary containers.
_PK_MAGIC = b"PK\x03\x04"
_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def _guard_ooxml_container(
    source_path: Path,
    source: str,
    target: str,
    legacy_label: str,
    canonical_label: str,
) -> None:
    """Reject anything that is not a real OOXML ZIP container, honestly."""
    try:
        with source_path.open("rb") as fh:
            header = fh.read(8)
    except OSError as exc:
        raise _unsupported(
            source,
            target,
            f"{source} to {target} conversion failed: could not read the source file ({exc}).",
        ) from exc

    if header.startswith(_OLE2_MAGIC):
        raise _unsupported(
            source,
            target,
            f"Legacy .{legacy_label} format (OLE2 compound document) is not supported "
            f"by this converter. Please save your document as .{canonical_label} and try again.",
        )
    if not header.startswith(_PK_MAGIC):
        raise _unsupported(
            source,
            target,
            f"The uploaded file is not a valid {canonical_label.upper()} document. "
            f"Please save it as a valid .{canonical_label} file and try again.",
        )


def _convert_docx_to_html(
    self, source_path: Path, target_format: str, working_root: Path
) -> Path:
    """mammoth semantic HTML, wrapped in a deterministic HTML5 document."""
    _guard_ooxml_container(source_path, "docx", "html", "doc", "docx")

    try:
        import mammoth
    except ImportError as exc:
        raise RuntimeError("mammoth is required for DOCX to HTML conversion.") from exc

    try:
        with source_path.open("rb") as fh:
            result = mammoth.convert_to_html(fh)
    except Exception as exc:
        raise _unsupported(
            "docx",
            "html",
            "DOCX to HTML conversion failed: the file is not a readable DOCX "
            f"document ({exc}).",
        ) from exc

    # mammoth returns "" for documents with no convertible content; the HTML5
    # wrapper below still guarantees a non-empty, valid, servable file.
    document = (
        "<!DOCTYPE html>\n"
        '<html lang="en">\n'
        "<head>\n"
        '<meta charset="utf-8">\n'
        f"<title>{_html.escape(source_path.stem)}</title>\n"
        "</head>\n"
        "<body>\n"
        f"{result.value}\n"
        "</body>\n</html>\n"
    )

    suffix = target_format.lower().lstrip(".")
    output_path = working_root / f"{source_path.stem}.{suffix}"
    output_path.write_text(document, encoding="utf-8")
    return output_path


def _render_slide_lines_to_pdf(lines: list[str], output_path: Path) -> Path:
    """Render text lines to a one-page reportlab PDF (engine JPG-path mirror)."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except ImportError as exc:
        raise RuntimeError(
            "reportlab is required for PPTX to PNG conversion."
        ) from exc

    output_path.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    margin = 40
    y = height - margin
    line_height = 14

    if not lines:
        lines = ["(slide contains non-text content)"]

    text_obj = c.beginText(margin, y)
    text_obj.setFont("Helvetica", 10)

    for line in lines:
        if y < margin + line_height:
            c.drawText(text_obj)
            c.showPage()
            y = height - margin
            text_obj = c.beginText(margin, y)
            text_obj.setFont("Helvetica", 10)

        text_obj.textLine(line[:120])
        y -= line_height

    c.drawText(text_obj)
    c.showPage()
    c.save()
    return output_path


def _convert_pptx_to_png(
    self, source_path: Path, target_format: str, working_root: Path
) -> Path:
    """First slide -> text lines -> reportlab PDF -> PyMuPDF PNG (engine profile)."""
    _guard_ooxml_container(source_path, "pptx", "png", "ppt", "pptx")

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX to PNG conversion.") from exc

    try:
        presentation = Presentation(str(source_path))
    except Exception as exc:
        raise _unsupported(
            "pptx",
            "png",
            "PPTX to PNG conversion failed: the file is not a readable PPTX "
            f"presentation ({exc}).",
        ) from exc

    if not presentation.slides:
        raise _unsupported(
            "pptx",
            "png",
            "PPTX to PNG conversion failed: the presentation contains no slides.",
        )

    # Exact extraction rules of DocumentEngine._convert_presentation_to_pdf,
    # restricted to the first slide (single-slide output, first-slide-only).
    slide = presentation.slides[0]
    lines: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(
                run.text or "" for run in paragraph.runs
            ).strip()
            if text:
                lines.append(text)
    if not lines:
        lines = ["(slide contains non-text content)"]

    pdf_path = working_root / f"{source_path.stem}.pdf"
    _render_slide_lines_to_pdf(lines, pdf_path)

    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for PPTX to PNG conversion.") from exc

    suffix = target_format.lower().lstrip(".")
    output_path = working_root / f"{source_path.stem}.{suffix}"
    pdf_doc = fitz.open(str(pdf_path))
    try:
        page = pdf_doc.load_page(0)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        pix.save(str(output_path))
    finally:
        pdf_doc.close()
        pdf_path.unlink(missing_ok=True)

    return output_path


DocxToHtmlPlugin = make_plugin_class(
    slug="docx-to-html",
    source_formats=["docx"],
    target_formats=["html"],
    engine_hook=_convert_docx_to_html,
    name="DOCX to HTML",
    description="Convert DOCX documents to clean, semantic HTML files.",
    category="document",
    engine="document",
    goal="document",
    use_case="Best for publishing Word documents as web-ready HTML pages.",
    priority=75,
    quality=85,
    compatibility=80,
    estimated_saving=8,
    badge="Office Conversion",
    icon="🌐",
    seo_title="DOCX to HTML Converter | Converigo",
    seo_description="Convert DOCX documents to HTML files quickly and easily.",
)

PptxToPngPlugin = make_plugin_class(
    slug="pptx-to-png",
    source_formats=["pptx"],
    target_formats=["png"],
    engine_hook=_convert_pptx_to_png,
    name="PPTX to PNG",
    description="Convert PowerPoint presentations to PNG images (first slide).",
    category="document",
    engine="document",
    goal="document",
    use_case="Best for turning PowerPoint slides into PNG images for previews and thumbnails.",
    priority=75,
    quality=85,
    compatibility=80,
    estimated_saving=8,
    badge="Office Conversion",
    icon="🖼️",
    seo_title="PPTX to PNG Converter | Converigo",
    seo_description="Convert PPTX presentations to PNG images quickly and easily.",
)