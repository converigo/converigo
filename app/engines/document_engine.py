import logging
from pathlib import Path
import tempfile

from PIL import Image

from app.engines.base_engine import BaseEngine

logger = logging.getLogger(__name__)


class DocumentEngine(BaseEngine):
    ENGINE_NAME = "document"

    SUPPORTED_FORMATS = [
        "pdf",
        "docx",
        "doc",
        "txt",
        "md",
        "xlsx",
        "xls",
        "pptx",
        "ppt",
        "odt",
    ]

    async def convert(
        self,
        source_path: Path,
        target_format: str,
    ) -> Path:
        target = target_format.lower().lstrip(".")
        if target == "doc":
            target = "docx"
        if target == "xls":
            target = "xlsx"
        if target == "ppt":
            target = "pptx"

        source_format = source_path.suffix.lower().lstrip(".")
        # Use configured OUTPUT_DIR to avoid environment-dependent CWD issues.
        # settings.OUTPUT_DIR is mapped from OUTPUT_DIR env var in app/core/settings.py
        from app.core.settings import settings

        output_dir = settings.OUTPUT_DIR / "document"
        output_dir.mkdir(parents=True, exist_ok=True)


        if source_format == "pdf":
            if target == "docx":
                return self._convert_pdf_to_docx(source_path, output_dir)

            if target == "xlsx":
                return self._convert_pdf_to_xlsx(source_path, output_dir)

            if target == "pptx":
                return self._convert_pdf_to_pptx(source_path, output_dir)

            if target == "odt":
                return self._convert_pdf_to_odt(source_path, output_dir)

            if target in {"jpg", "jpeg"}:
                return self._convert_pdf_to_jpg(source_path, output_dir)

            raise ValueError(
                f"Unsupported target format for document engine: {target}"
            )

        if target != "pdf":
            raise ValueError(
                f"Unsupported target format for document engine: {target}"
            )

        if source_format in {"xlsx", "xls", "csv"}:
            return self._convert_spreadsheet_to_pdf(source_path, output_dir)

        if source_format in {"pptx", "ppt"}:
            return self._convert_presentation_to_pdf(source_path, output_dir)

        if source_format == "odt":
            return self._convert_odt_to_pdf(source_path, output_dir)

        raise ValueError(
            f"Unsupported source format for document engine: {source_path.suffix}"
        )

    def _render_text_lines_to_pdf(
        self,
        lines: list[str],
        output_path: Path,
    ) -> Path:
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
        except ImportError as exc:
            raise RuntimeError(
                "reportlab is required for office document to PDF conversion."
            ) from exc

        output_path.parent.mkdir(parents=True, exist_ok=True)
        c = canvas.Canvas(str(output_path), pagesize=letter)
        width, height = letter
        margin = 40
        y = height - margin
        line_height = 14

        if not lines:
            lines = ["(no text extracted from document)"]

        text_obj = c.beginText(margin, y)
        text_obj.setFont("Helvetica", 10)

        for line in lines:
            if y < margin + line_height:
                c.drawText(text_obj)
                c.showPage()
                y = height - margin
                text_obj = c.beginText(margin, y)
                text_obj.setFont("Helvetica", 10)

            text_obj.textLine(line[:120])
            y -= line_height

        c.drawText(text_obj)
        c.showPage()
        c.save()
        return output_path

    def _convert_spreadsheet_to_pdf(
        self,
        source_path: Path,
        output_dir: Path,
    ) -> Path:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise RuntimeError(
                "openpyxl is required for XLSX/XLS to PDF conversion."
            ) from exc

        output_path = output_dir / f"{source_path.stem}.pdf"
        lines: list[str] = []

        if source_path.suffix.lower() in {".xlsx", ".xls"}:
            workbook = load_workbook(str(source_path), data_only=True, read_only=True)
            for sheet in workbook.worksheets:
                lines.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    if row and any(cell is not None for cell in row):
                        row_text = " | ".join(
                            "" if cell is None else str(cell)
                            for cell in row
                        )
                        lines.append(row_text)
                lines.append("")
        else:
            import csv

            lines.append(f"CSV: {source_path.name}")
            with source_path.open(newline="", encoding="utf-8") as csv_file:
                reader = csv.reader(csv_file)
                for row in reader:
                    lines.append(" | ".join(row))

        return self._render_text_lines_to_pdf(lines, output_path)

    def _convert_presentation_to_pdf(
        self,
        source_path: Path,
        output_dir: Path,
    ) -> Path:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is required for PPTX/PPT to PDF conversion."
            ) from exc

        output_path = output_dir / f"{source_path.stem}.pdf"
        presentation = Presentation(str(source_path))
        lines: list[str] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"Slide: {slide_index}")
            found_text = False
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(
                        run.text or "" for run in paragraph.runs
                    ).strip()
                    if text:
                        lines.append(text)
                        found_text = True
            if not found_text:
                lines.append("(slide contains non-text content)")
            lines.append("")

        return self._render_text_lines_to_pdf(lines, output_path)

    def _convert_odt_to_pdf(
        self,
        source_path: Path,
        output_dir: Path,
    ) -> Path:
        try:
            from odf.opendocument import load
            from odf import text as odf_text
        except ImportError as exc:
            raise RuntimeError(
                "odfpy is required for ODT to PDF conversion."
            ) from exc

        output_path = output_dir / f"{source_path.stem}.pdf"
        document = load(str(source_path))
        lines: list[str] = []

        for paragraph in document.getElementsByType(odf_text.P):
            text = self._extract_odt_text(paragraph).strip()
            if text:
                lines.append(text)

        if not lines:
            lines.append("(no extractable ODT text found)")

        return self._render_text_lines_to_pdf(lines, output_path)

    def _extract_odt_text(self, node) -> str:
        text_parts: list[str] = []

        if hasattr(node, "childNodes"):
            for child in node.childNodes:
                text_parts.append(self._extract_odt_text(child))
        else:
            if getattr(node, "data", None):
                text_parts.append(node.data)

        return "".join(text_parts)

    def _validate_pdf_input(self, source_path: Path) -> None:
        """
        Validate PDF input before conversion.
        
        Raises:
            PDFEmptyError: If PDF has no pages.
            PDFPasswordProtectedError: If PDF is password protected.
            PDFValidationError: For other unexpected validation failures.
        """
        # Import exceptions locally to avoid circular imports
        from app.services.conversion_service import (
            PDFEmptyError,
            PDFPasswordProtectedError,
            PDFValidationError,
        )
        
        try:
            import pdfplumber
        except ImportError as exc:
            raise PDFValidationError("pdfplumber is required for PDF validation.") from exc

        try:
            with pdfplumber.open(str(source_path)) as pdf:
                # Check if PDF has any pages
                if len(pdf.pages) == 0:
                    raise PDFEmptyError()
                
                logger.info("PDF validation passed: %d pages found", len(pdf.pages))
                
        except PDFEmptyError:
            # Re-raise our custom validation errors
            raise
        except PDFPasswordProtectedError:
            # Re-raise our custom validation errors
            raise
        except Exception as exc:
            # Check if it's a permission/password error
            exc_str = str(exc).lower()
            exc_type = type(exc).__name__.lower()
            
            if "permission" in exc_str or "password" in exc_str or "encrypted" in exc_str:
                raise PDFPasswordProtectedError() from exc
            
            # For other unexpected errors during validation, log and raise a generic validation error
            logger.exception("Unexpected error during PDF validation: %s", exc)
            raise PDFValidationError(
                f"PDF validation failed: {type(exc).__name__}"
            ) from exc

    def _convert_pdf_to_docx(self, source_path: Path, output_dir: Path) -> Path:
        # Validate PDF input before conversion
        self._validate_pdf_input(source_path)
        
        try:
            from pdf2docx import Converter
        except ImportError as exc:
            raise RuntimeError("pdf2docx is required for PDF to DOCX conversion.") from exc

        output_path = output_dir / f"{source_path.stem}.docx"

        try:
            converter = Converter(str(source_path))
            converter.convert(str(output_path), start=0, end=None)
            converter.close()
        except Exception as exc:
            # Handle the specific pdf2docx failure where no pages are parsed
            exc_str = str(exc)
            try:
                from pdf2docx import converter as _pdf2conv
                is_conv_exc = isinstance(exc, getattr(_pdf2conv, 'ConversionException', Exception))
            except Exception:
                is_conv_exc = False

            if is_conv_exc or 'No parsed pages' in exc_str:
                try:
                    import pdfplumber
                    from docx import Document
                except ImportError as imp_exc:
                    raise RuntimeError(
                        "Fallback conversion requires pdfplumber and python-docx"
                    ) from imp_exc

                # Attempt a minimal fallback: extract text and write a simple DOCX
                with pdfplumber.open(str(source_path)) as pdf:
                    pages_text: list[str] = []
                    for page in pdf.pages:
                        txt = page.extract_text() or ""
                        if txt.strip():
                            pages_text.append(txt)

                doc = Document()
                if pages_text:
                    for page_txt in pages_text:
                        for line in page_txt.splitlines():
                            doc.add_paragraph(line)
                else:
                    doc.add_paragraph("(no extractable text found in PDF)")

                doc.save(str(output_path))
            else:
                raise RuntimeError(f"PDF to DOCX conversion failed: {exc}") from exc

        if not output_path.exists():
            raise RuntimeError("PDF to DOCX conversion did not produce output.")

        return output_path

    def _convert_pdf_to_xlsx(self, source_path: Path, output_dir: Path) -> Path:
        try:
            import pdfplumber
            from openpyxl import Workbook
        except ImportError as exc:
            raise RuntimeError(
                "pdfplumber and openpyxl are required for PDF to XLSX conversion."
            ) from exc

        output_path = output_dir / f"{source_path.stem}.xlsx"
        logger.info("PDF input path: %s", source_path)

        try:
            with pdfplumber.open(str(source_path)) as pdf:
                workbook = Workbook()
                workbook.remove(workbook.active)
                found_table = False

                for page_index, page in enumerate(pdf.pages, start=1):
                    tables = page.extract_tables() or []
                    if tables:
                        found_table = True
                        sheet = workbook.create_sheet(title=f"page_{page_index}")
                        for row_index, row in enumerate(tables[0], start=1):
                            for col_index, cell in enumerate(row, start=1):
                                sheet.cell(row=row_index, column=col_index, value=cell)

                if not found_table:
                    sheet = workbook.create_sheet(title="page_1")
                    for page_index, page in enumerate(pdf.pages, start=1):
                        text = page.extract_text() or ""
                        for line_index, line in enumerate(text.splitlines(), start=1):
                            sheet.cell(row=line_index, column=1, value=line)

                workbook.save(str(output_path))
                logger.info("XLSX output path: %s", output_path)
        except Exception as exc:
            logger.exception("PDF to XLSX conversion failed")
            raise RuntimeError(f"PDF to XLSX conversion failed: {type(exc).__name__}: {exc}") from exc

        if not output_path.exists():
            raise RuntimeError("PDF to XLSX conversion did not produce output.")

        return output_path

    def _convert_pdf_to_pptx(self, source_path: Path, output_dir: Path) -> Path:
        try:
            import fitz
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:
            raise RuntimeError(
                "PyMuPDF and python-pptx are required for PDF to PPTX conversion."
            ) from exc

        output_path = output_dir / f"{source_path.stem}.pptx"

        try:
            presentation = Presentation()
            blank_slide_layout = presentation.slide_layouts[6]

            with tempfile.TemporaryDirectory(prefix="pdf2pptx_", dir=str(output_dir)) as temp_dir:
                doc = fitz.open(str(source_path))
                try:
                    if doc.page_count == 0:
                        raise RuntimeError("No pages were extracted from the PDF.")

                    for page_index in range(doc.page_count):
                        page = doc.load_page(page_index)
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                        image_path = Path(temp_dir) / f"page_{page_index + 1:02d}.png"
                        pix.save(str(image_path))

                        slide = presentation.slides.add_slide(blank_slide_layout)
                        slide.shapes.add_picture(
                            str(image_path),
                            Inches(0),
                            Inches(0),
                            width=Inches(10),
                            height=Inches(7.5),
                        )

                    presentation.save(str(output_path))
                finally:
                    doc.close()
        except Exception as exc:
            raise RuntimeError(f"PDF to PPTX conversion failed: {exc}") from exc

        if not output_path.exists():
            raise RuntimeError("PDF to PPTX conversion did not produce output.")

        return output_path

    def _convert_pdf_to_odt(self, source_path: Path, output_dir: Path) -> Path:
        try:
            import pdfplumber
            from odf.opendocument import OpenDocumentText
            from odf.text import P
        except ImportError as exc:
            raise RuntimeError(
                "pdfplumber and odfpy are required for PDF to ODT conversion."
            ) from exc

        output_path = output_dir / f"{source_path.stem}.odt"

        try:
            doc = OpenDocumentText()
            with pdfplumber.open(str(source_path)) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    for line in text.splitlines():
                        paragraph = P(text=line)
                        doc.text.addElement(paragraph)

            doc.save(str(output_path))
        except Exception as exc:
            raise RuntimeError(f"PDF to ODT conversion failed: {exc}") from exc

        if not output_path.exists():
            raise RuntimeError("PDF to ODT conversion did not produce output.")

        return output_path

    def _convert_pdf_to_jpg(self, source_path: Path, output_dir: Path) -> Path:
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError("PyMuPDF is required for PDF to JPG conversion.") from exc

        output_path = output_dir / f"{source_path.stem}_page_01.jpg"

        try:
            doc = fitz.open(str(source_path))
            try:
                if doc.page_count == 0:
                    raise RuntimeError("No pages were extracted from the PDF.")

                page = doc.load_page(0)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                pix.pil_save(str(output_path), format="JPEG", quality=95)
                return output_path
            finally:
                doc.close()
        except Exception as exc:
            raise RuntimeError(f"PDF to JPG conversion failed: {exc}") from exc
